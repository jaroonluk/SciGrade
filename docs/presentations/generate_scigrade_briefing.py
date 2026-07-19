# -*- coding: utf-8 -*-
"""Generate SciGrade executive PowerPoint + A4 one-pager."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
    KeepTogether,
)

OUT_DIR = Path(__file__).resolve().parent

# Brand colors (SciGrade)
BROWN = RGBColor(0x5C, 0x2E, 0x1F)
AMBER = RGBColor(0x8B, 0x45, 0x13)
CREAM = RGBColor(0xFD, 0xF6, 0xF0)
DARK = RGBColor(0x2D, 0x2A, 0x26)
MUTED = RGBColor(0x7A, 0x4A, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_LINE = RGBColor(0xE8, 0xD5, 0xC4)

RL_BROWN = colors.HexColor("#5C2E1F")
RL_AMBER = colors.HexColor("#8B4513")
RL_CREAM = colors.HexColor("#FDF6F0")
RL_MUTED = colors.HexColor("#7A4A3A")
RL_LINE = colors.HexColor("#E8D5C4")
RL_DARK = colors.HexColor("#2D2A26")


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, prs, color=CREAM):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_top_bar(slide, prs):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()


def add_footer(slide, prs, page, total=8):
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.0), Inches(8.5), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"SciGrade — คณะวิทยาศาสตร์  |  เอกสารนำเสนอผู้บริหาร  |  {page}/{total}"
    set_run(run, 10, False, MUTED)
    p.alignment = PP_ALIGN.LEFT


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, BROWN)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(5.8), prs.slide_width, Inches(1.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.5), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SciGrade"
    set_run(r, 48, True, WHITE)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "ระบบรายงานผลการสอบ"
    set_run(r2, 28, True, WHITE)

    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "โครงสร้างระบบและภาพรวมการใช้งาน"
    set_run(r3, 18, False, CREAM)

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(8.5), Inches(1.0))
    ftf = foot.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "คณะวิทยาศาสตร์\nเอกสารสรุปสำหรับนำเสนอผู้บริหาร"
    set_run(fr, 14, False, WHITE)


def content_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, 26, True, BROWN)


def add_bullets(slide, left, top, width, height, items, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = item
        set_run(r, size, False, DARK)
        p.space_after = Pt(8)


def card(slide, left, top, width, height, title, body_lines, title_size=14, body_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_LINE
    shape.adjustments[0] = 0.08

    tbox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.4)
    )
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, title_size, True, AMBER)

    bbox = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.6),
    )
    btf = bbox.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        r = p.add_run()
        r.text = line
        set_run(r, body_size, False, DARK)
        p.space_after = Pt(4)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "1. ภาพรวมระบบ")
    add_bullets(
        slide,
        Inches(0.6),
        Inches(1.1),
        Inches(8.8),
        Inches(5.5),
        [
            "SciGrade เป็นระบบเว็บสำหรับจัดการรายงานผลการสอบไล่ ของคณะวิทยาศาสตร์",
            "ครอบคลุมตั้งแต่กรอก/อัปโหลดโดยอาจารย์ → อนุมัติสาขา → อนุมัติคณะ → พิมพ์เอกสาร",
            "ดึงรายวิชาจากฐานข้อมูล REG ของมหาวิทยาลัยมาใช้ติดตามสถานะการส่งผลสอบ",
            "เข้าสู่ระบบด้วย Google Account (เฉพาะบุคลากรที่มีในฐานข้อมูลคณะ)",
            "เทคโนโลยี: Laravel 12, MySQL (SciGrad + REG), Docker",
            "จุดสำคัญ: อ่านข้อมูลจาก REG เท่านั้น — ไม่เขียนกลับเข้า REG",
        ],
        size=16,
    )
    add_footer(slide, prs, 2)


def slide_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "2. ผู้ใช้งานและสิทธิ์")

    cards = [
        ("อาจารย์ผู้สอน", ["สร้าง/อัปโหลดรายงาน", "แก้ไขเมื่อถูกส่งกลับ", "พิมพ์รายงานของตนเอง"]),
        ("Admin สาขา", ["อนุมัติระดับสาขา", "ติดตามสถานะส่งผล", "ส่งเอกสารให้คณะ", "พิมพ์รายงานสาขา"]),
        ("Admin กลาง", ["รับเอกสารจากสาขา", "อนุมัติระดับคณะ", "จัดการ REG / ตั้งค่า", "จัดการสิทธิ์ผู้ใช้"]),
        ("Super Admin", ["สิทธิ์ Admin กลางทั้งหมด", "กำหนดรูปแบบรหัสวิชา", "สลับบัญชีใช้งาน"]),
    ]
    for i, (title, lines) in enumerate(cards):
        col = i % 2
        row = i // 2
        card(
            slide,
            Inches(0.5 + col * 4.6),
            Inches(1.1 + row * 2.7),
            Inches(4.35),
            Inches(2.45),
            title,
            lines,
            16,
            13,
        )
    add_footer(slide, prs, 3)


def slide_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "3. โครงสร้างโมดูลหลัก")

    card(
        slide,
        Inches(0.4),
        Inches(1.05),
        Inches(3.0),
        Inches(5.5),
        "ฝั่งอาจารย์",
        [
            "• สร้างรายงานผลสอบ",
            "• อัปโหลดจาก PDF ทะเบียน",
            "• รายงานของฉัน",
            "• พิมพ์รายงานรายวิชา",
        ],
        15,
        13,
    )
    card(
        slide,
        Inches(3.5),
        Inches(1.05),
        Inches(3.0),
        Inches(5.5),
        "ฝั่งสาขา",
        [
            "• ตรวจสอบรายวิชา",
            "• สถานะการส่ง REG",
            "• ส่งเอกสารให้คณะ",
            "• พิมพ์รายงานสาขา",
            "  (PDF / Word)",
        ],
        15,
        13,
    )
    card(
        slide,
        Inches(6.6),
        Inches(1.05),
        Inches(3.0),
        Inches(5.5),
        "ฝั่งคณะ / Admin กลาง",
        [
            "• รับเอกสารจากสาขา",
            "• อนุมัติระดับคณะ",
            "• รายงานทุกสาขา",
            "• Download / จัดการ REG",
            "• ตั้งค่าภาค / หลักสูตร",
            "• จัดการสิทธิ์ผู้ใช้",
            "• จัดกลุ่มรายวิชา",
        ],
        15,
        12,
    )
    add_footer(slide, prs, 4)


def slide_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "4. กระบวนการอนุมัติผลสอบ")

    steps = [
        ("1", "อาจารย์", "สร้างหรืออัปโหลด\nรายงานผลสอบ"),
        ("2", "สาขา", "ตรวจสอบและ\nอนุมัติระดับสาขา"),
        ("3", "สาขา→คณะ", "ส่งชุดเอกสาร\nคณะรับเอกสาร"),
        ("4", "คณะ", "อนุมัติระดับคณะ\nพิมพ์/ส่งออก"),
    ]
    for i, (num, role, desc) in enumerate(steps):
        x = Inches(0.45 + i * 2.4)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.25), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = AMBER
        circ.line.fill.background()
        nbox = slide.shapes.add_textbox(x + Inches(0.7), Inches(1.32), Inches(0.55), Inches(0.45))
        np = nbox.text_frame.paragraphs[0]
        nr = np.add_run()
        nr.text = num
        set_run(nr, 16, True, WHITE)
        np.alignment = PP_ALIGN.CENTER

        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + Inches(1.85),
                Inches(1.4),
                Inches(0.4),
                Inches(0.25),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_LINE
            arrow.line.fill.background()

        card(slide, x, Inches(2.1), Inches(2.2), Inches(2.2), role, [desc], 15, 12)

    # Status table
    status_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.4))
    st = status_title.text_frame.paragraphs[0]
    sr = st.add_run()
    sr.text = "สถานะรายงานผลสอบ"
    set_run(sr, 16, True, BROWN)

    statuses = [
        ("0", "บันทึกแล้ว / รอสาขา"),
        ("1", "ผ่านระดับสาขา"),
        ("2", "ผ่านระดับคณะ"),
        ("-1", "ส่งกลับแก้ไข"),
    ]
    for i, (code, label) in enumerate(statuses):
        card(
            slide,
            Inches(0.45 + i * 2.35),
            Inches(5.0),
            Inches(2.2),
            Inches(1.5),
            f"สถานะ {code}",
            [label],
            13,
            12,
        )
    add_footer(slide, prs, 5)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "5. สถาปัตยกรรมและการเชื่อมต่อ")

    card(
        slide,
        Inches(0.5),
        Inches(1.1),
        Inches(4.4),
        Inches(3.2),
        "สถาปัตยกรรม",
        [
            "• ผู้ใช้เข้าผ่านเว็บ Browser",
            "• แอป SciGrade (Laravel + Nginx)",
            "• ฐาน SciGrad: รายงาน/สิทธิ์/เอกสาร",
            "• ฐาน REG: รายวิชา/กลุ่มเรียน (อ่านอย่างเดียว)",
            "• ยืนยันตัวตนด้วย Google OAuth",
        ],
        15,
        13,
    )
    card(
        slide,
        Inches(5.1),
        Inches(1.1),
        Inches(4.4),
        Inches(3.2),
        "ระบบภายนอก",
        [
            "• Google Account — เข้าสู่ระบบ",
            "• ฐานข้อมูล REG — ดึงรายวิชา",
            "• ไฟล์ PDF จากทะเบียน — อัปโหลดแล้วระบบอ่านข้อมูล",
        ],
        15,
        13,
    )
    card(
        slide,
        Inches(0.5),
        Inches(4.5),
        Inches(9.0),
        Inches(2.1),
        "เอกสารที่ระบบผลิตได้",
        [
            "พิมพ์รายงานรายวิชาบนเบราว์เซอร์  •  ส่งออก PDF  •  ส่งออก Word (.docx)  •  ไฟล์แนบจากอาจารย์และชุดเอกสารสาขา",
        ],
        15,
        14,
    )
    add_footer(slide, prs, 6)


def slide_benefits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)
    content_title(slide, "6. ประโยชน์ต่อหน่วยงาน")
    items = [
        ("ลดงานเอกสารซ้ำ", "ติดตามสถานะออนไลน์แทนการไล่ถามหลายช่องทาง"),
        ("ชัดเจนเรื่องสิทธิ์", "แยกบทบาทอาจารย์ / สาขา / คณะ อย่างเป็นระบบ"),
        ("ตรวจสอบย้อนกลับได้", "มีบันทึกการอนุมัติและสถานะในแต่ละขั้น"),
        ("สอดคล้องกับ REG", "ใช้รายวิชาจากระบบทะเบียนเป็นฐานติดตาม"),
        ("พร้อมส่งออกเอกสาร", "PDF / Word สำหรับการประชุมและการจัดเก็บ"),
    ]
    for i, (t, d) in enumerate(items):
        y = Inches(1.05 + i * 1.05)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y, Inches(8.9), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_LINE
        shape.adjustments[0] = 0.1

        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.75), y + Inches(0.2), Inches(0.5), Inches(0.5)
        )
        num.fill.solid()
        num.fill.fore_color.rgb = AMBER
        num.line.fill.background()
        nbox = slide.shapes.add_textbox(Inches(0.75), y + Inches(0.28), Inches(0.5), Inches(0.4))
        np = nbox.text_frame.paragraphs[0]
        nr = np.add_run()
        nr.text = str(i + 1)
        set_run(nr, 14, True, WHITE)
        np.alignment = PP_ALIGN.CENTER

        tbox = slide.shapes.add_textbox(Inches(1.45), y + Inches(0.15), Inches(7.7), Inches(0.65))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = t
        set_run(r, 15, True, BROWN)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = d
        set_run(r2, 12, False, DARK)

    add_footer(slide, prs, 7)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, BROWN)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(3.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "สรุป"
    set_run(r, 36, True, WHITE)

    lines = [
        "SciGrade เชื่อมโยงงานรายงานผลสอบทั้งสายงาน",
        "อาจารย์ → สาขา → คณะ อย่างเป็นขั้นตอน ตรวจสอบได้",
        "และใช้ข้อมูล REG เป็นฐานในการติดตาม",
    ]
    for line in lines:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = line
        set_run(r2, 18, False, CREAM)
        p2.space_before = Pt(6)

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.6))
    fp = foot.text_frame.paragraphs[0]
    fr = fp.add_run()
    fr.text = "คณะวิทยาศาสตร์  •  SciGrade"
    set_run(fr, 14, False, WHITE)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    slide_overview(prs)
    slide_roles(prs)
    slide_modules(prs)
    slide_workflow(prs)
    slide_architecture(prs)
    slide_benefits(prs)
    slide_closing(prs)

    out = OUT_DIR / "SciGrade_System_Overview_Exec.pptx"
    prs.save(str(out))
    return out


def build_a4_pdf():
    """One-page A4 executive handout."""
    out = OUT_DIR / "SciGrade_OnePage_A4.pdf"

    # Try Thai-capable fonts on Windows
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    font_candidates = [
        (r"C:\Windows\Fonts\THSarabunNew.ttf", r"C:\Windows\Fonts\THSarabunNew Bold.ttf"),
        (r"C:\Windows\Fonts\thsarabunnew.ttf", r"C:\Windows\Fonts\thsarabunnew_bold.ttf"),
        (r"C:\Windows\Fonts\cordia.ttf", r"C:\Windows\Fonts\cordiab.ttf"),
        (r"C:\Windows\Fonts\tahoma.ttf", r"C:\Windows\Fonts\tahomabd.ttf"),
        (r"C:\Windows\Fonts\arial.ttf", r"C:\Windows\Fonts\arialbd.ttf"),
    ]
    font_name = "Helvetica"
    font_bold = "Helvetica-Bold"
    for regular, bold in font_candidates:
        if Path(regular).exists():
            try:
                pdfmetrics.registerFont(TTFont("ThaiR", regular))
                bold_path = bold if Path(bold).exists() else regular
                pdfmetrics.registerFont(TTFont("ThaiB", bold_path))
                font_name = "ThaiR"
                font_bold = "ThaiB"
                break
            except Exception:
                continue

    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=12 * mm,
        rightMargin=12 * mm,
        topMargin=10 * mm,
        bottomMargin=8 * mm,
    )

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="H1",
            fontName=font_bold,
            fontSize=16 if font_name != "Helvetica" else 14,
            leading=20,
            textColor=RL_BROWN,
            alignment=TA_CENTER,
            spaceAfter=2,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Sub",
            fontName=font_name,
            fontSize=10 if font_name != "Helvetica" else 9,
            leading=13,
            textColor=RL_MUTED,
            alignment=TA_CENTER,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Sec",
            fontName=font_bold,
            fontSize=11 if font_name != "Helvetica" else 10,
            leading=14,
            textColor=RL_AMBER,
            spaceBefore=4,
            spaceAfter=2,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Body",
            fontName=font_name,
            fontSize=9 if font_name != "Helvetica" else 8,
            leading=12,
            textColor=RL_DARK,
            spaceAfter=1,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Cell",
            fontName=font_name,
            fontSize=8 if font_name != "Helvetica" else 7.5,
            leading=10,
            textColor=RL_DARK,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CellB",
            fontName=font_bold,
            fontSize=8 if font_name != "Helvetica" else 7.5,
            leading=10,
            textColor=RL_BROWN,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Tiny",
            fontName=font_name,
            fontSize=7.5 if font_name != "Helvetica" else 7,
            leading=9,
            textColor=RL_MUTED,
            alignment=TA_CENTER,
        )
    )

    story = []
    story.append(Paragraph("SciGrade — ระบบรายงานผลการสอบ", styles["H1"]))
    story.append(
        Paragraph(
            "คณะวิทยาศาสตร์  •  สรุปโครงสร้างระบบสำหรับผู้บริหาร (1 หน้า)",
            styles["Sub"],
        )
    )

    story.append(Paragraph("1. ภาพรวม", styles["Sec"]))
    story.append(
        Paragraph(
            "ระบบเว็บจัดการรายงานผลการสอบไล่ ตั้งแต่กรอก/อัปโหลดโดยอาจารย์ "
            "→ อนุมัติระดับสาขา → ส่งเอกสารและอนุมัติระดับคณะ → พิมพ์/ส่งออกเอกสาร "
            "พร้อมดึงรายวิชาจากฐาน REG (อ่านอย่างเดียว ไม่เขียนกลับ) เข้าสู่ระบบด้วย Google Account",
            styles["Body"],
        )
    )

    story.append(Paragraph("2. ผู้ใช้งาน", styles["Sec"]))
    role_data = [
        [
            Paragraph("<b>บทบาท</b>", styles["CellB"]),
            Paragraph("<b>หน้าที่หลัก</b>", styles["CellB"]),
        ],
        [
            Paragraph("อาจารย์ผู้สอน", styles["CellB"]),
            Paragraph("สร้าง/อัปโหลดรายงาน แก้ไขเมื่อถูกส่งกลับ พิมพ์รายงานของตนเอง", styles["Cell"]),
        ],
        [
            Paragraph("Admin สาขา", styles["CellB"]),
            Paragraph("อนุมัติสาขา ติดตามสถานะ ส่งเอกสารให้คณะ พิมพ์รายงานสาขา", styles["Cell"]),
        ],
        [
            Paragraph("Admin กลาง", styles["CellB"]),
            Paragraph("รับเอกสาร อนุมัติคณะ จัดการ REG/ตั้งค่า จัดการสิทธิ์ผู้ใช้", styles["Cell"]),
        ],
        [
            Paragraph("Super Admin", styles["CellB"]),
            Paragraph("สิทธิ์ Admin กลางทั้งหมด + กำหนดรูปแบบรหัสวิชา + สลับบัญชี", styles["Cell"]),
        ],
    ]
    t1 = Table(role_data, colWidths=[32 * mm, 148 * mm])
    t1.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), RL_CREAM),
                ("GRID", (0, 0), (-1, -1), 0.4, RL_LINE),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    story.append(t1)

    story.append(Paragraph("3. โมดูลหลัก", styles["Sec"]))
    mod_data = [
        [
            Paragraph("<b>อาจารย์</b>", styles["CellB"]),
            Paragraph("<b>สาขา</b>", styles["CellB"]),
            Paragraph("<b>คณะ / Admin กลาง</b>", styles["CellB"]),
        ],
        [
            Paragraph("สร้างรายงาน<br/>อัปโหลด PDF ทะเบียน<br/>รายงานของฉัน", styles["Cell"]),
            Paragraph("ตรวจสอบรายวิชา<br/>สถานะส่ง REG<br/>ส่งเอกสารให้คณะ<br/>พิมพ์รายงานสาขา", styles["Cell"]),
            Paragraph(
                "รับเอกสาร อนุมัติคณะ<br/>รายงานทุกสาขา<br/>Download/จัดการ REG<br/>ตั้งค่าภาค/หลักสูตร/สิทธิ์",
                styles["Cell"],
            ),
        ],
    ]
    t2 = Table(mod_data, colWidths=[60 * mm, 60 * mm, 60 * mm])
    t2.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), RL_CREAM),
                ("GRID", (0, 0), (-1, -1), 0.4, RL_LINE),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    story.append(t2)

    story.append(Paragraph("4. Workflow การอนุมัติ", styles["Sec"]))
    flow = (
        "① อาจารย์สร้าง/อัปโหลด (สถานะ 0)  →  "
        "② สาขาอนุมัติ (สถานะ 1)  →  "
        "③ ส่งชุดเอกสารให้คณะรับ  →  "
        "④ คณะอนุมัติ (สถานะ 2)  →  พิมพ์ PDF/Word<br/>"
        "สถานะ -1 = ส่งกลับให้อาจารย์แก้ไข แล้วส่งการแก้ไขกลับเข้าสู่สายงาน"
    )
    story.append(Paragraph(flow, styles["Body"]))

    story.append(Paragraph("5. สถาปัตยกรรม / การเชื่อมต่อ", styles["Sec"]))
    story.append(
        Paragraph(
            "Browser → SciGrade (Laravel) → ฐาน SciGrad (รายงาน/สิทธิ์/เอกสาร) "
            "+ ฐาน REG (รายวิชา อ่านอย่างเดียว) + Google OAuth "
            "• ส่งออก: พิมพ์บนเว็บ / PDF / Word",
            styles["Body"],
        )
    )

    story.append(Paragraph("6. ประโยชน์", styles["Sec"]))
    story.append(
        Paragraph(
            "ลดงานเอกสารซ้ำ  •  ชัดเจนเรื่องสิทธิ์และความรับผิดชอบ  •  ตรวจสอบย้อนกลับได้  •  "
            "สอดคล้องข้อมูล REG  •  พร้อมส่งออกเอกสารสำหรับการประชุม",
            styles["Body"],
        )
    )

    story.append(Spacer(1, 4 * mm))
    story.append(
        Paragraph(
            "SciGrade  •  คณะวิทยาศาสตร์  •  เอกสารสรุปสำหรับผู้บริหาร",
            styles["Tiny"],
        )
    )

    doc.build(story)
    return out


if __name__ == "__main__":
    pptx_path = build_pptx()
    pdf_path = build_a4_pdf()
    print("PPTX:", pptx_path)
    print("PDF:", pdf_path)
